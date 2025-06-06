"""Utility module for extracting page counts from different file types.

This module provides functions to get page counts from various document types:
* PDF files using PyMuPDF
* Office documents (DOCX, PPTX) using python-docx and python-pptx
* Image files (JPG, TIFF, etc.) - always returns 1
"""

import io
import logging
from typing import Optional, Union

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

logger = logging.getLogger("page_count_util")


def get_page_count(file_input: Union[str, bytes], file_type: str) -> Optional[int]:
    """Get the page count from a document.

    :param file_input: Either a file path (str) or file content (bytes)
    :type file_input: Union[str, bytes]
    :param file_type: The type of the file (pdf, docx, pptx, jpg, tiff, etc.)
    :type file_type: str
    :return: Number of pages in the document, or None if counting fails
    :rtype: Optional[int]
    :raises ValueError: If file_type is not supported
    """
    file_type = file_type.lower()

    try:
        if file_type == "pdf":
            return _get_pdf_page_count(file_input)
        elif file_type == "docx":
            return _get_docx_page_count(file_input)
        elif file_type == "pptx":
            return _get_pptx_page_count(file_input)
        elif file_type == "xlsx":
            return 1  # excel always returns 1 page
        elif file_type in ["eml", "msg"]:
            return 1  # email always returns 1 page
        elif file_type in ["jpg", "jpeg", "png", "gif", "bmp", "tiff", "webp"]:
            return 1  # Image files always have 1 page
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    except Exception as e:
        logger.error(f"Error getting page count for {file_type} file: {str(e)}")
        return None


def _get_pdf_page_count(file_input: Union[str, bytes]) -> int:
    """Get page count from a PDF file.

    :param file_input: Either a file path or PDF content in bytes
    :type file_input: Union[str, bytes]
    :return: Number of pages in the PDF
    :rtype: int
    """
    if isinstance(file_input, str):
        # Input is a file path
        doc = fitz.open(file_input)
    else:
        # Input is bytes
        doc = fitz.open(stream=file_input, filetype="pdf")

    try:
        return doc.page_count
    finally:
        doc.close()


def _get_docx_page_count(file_input: Union[str, bytes]) -> int:
    """Get page count from a DOCX file.

    :param file_input: Either a file path or DOCX content in bytes
    :type file_input: Union[str, bytes]
    :return: Approximate number of pages in the document
    :rtype: int
    """
    if isinstance(file_input, str):
        doc = Document(file_input)
    else:
        doc = Document(io.BytesIO(file_input))

    # This is an approximation as actual page count depends on formatting
    return len(doc.sections)


def _get_pptx_page_count(file_input: Union[str, bytes]) -> int:
    """Get page count from a PPTX file.

    :param file_input: Either a file path or PPTX content in bytes
    :type file_input: Union[str, bytes]
    :return: Number of slides in the presentation
    :rtype: int
    """
    if isinstance(file_input, str):
        prs = Presentation(file_input)
    else:
        prs = Presentation(io.BytesIO(file_input))

    return len(prs.slides)
